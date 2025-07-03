import json
from datetime import date
from pathlib import Path
from typing import Any, Dict, Iterable, Union

try:
    from openpyxl import Workbook
    from openpyxl.utils import get_column_letter
except ImportError:  # pragma: no cover - openpyxl may not be installed
    Workbook = None
    get_column_letter = None

try:
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches, Pt
except ImportError:  # pragma: no cover - python-pptx may not be installed
    Presentation = None
    Inches = Pt = None

JsonLike = Union[str, Path, Dict[str, Any]]


def _load_json(data: JsonLike) -> Dict[str, Any]:
    """Load JSON from a path, JSON string, or dictionary."""
    if isinstance(data, dict):
        return data
    if isinstance(data, Path) or (isinstance(data, str) and Path(data).exists()):
        text = Path(data).read_text(encoding="utf-8")
        return json.loads(text)
    if isinstance(data, str):
        return json.loads(data)
    raise TypeError("Unsupported data type for result JSON")


def _slugify(value: str) -> str:
    return "".join(c if c.isalnum() else "_" for c in value)


def generate_excel_report(result_data: JsonLike, metadata: Dict[str, Any], output_dir: Union[str, Path] = "Reports") -> str:
    """Generate an Excel workbook summarizing valuation results."""
    if Workbook is None:
        raise ImportError("openpyxl is required to generate Excel reports")

    results = _load_json(result_data)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    case_name = metadata.get("case_name", "case")
    case_slug = _slugify(case_name) or "case"
    workbook_path = output_dir / f"report_{case_slug}.xlsx"

    wb = Workbook()
    default = wb.active
    wb.remove(default)

    for key, value in results.items():
        ws = wb.create_sheet(title=str(key)[:31])
        if isinstance(value, list):
            if value and isinstance(value[0], dict):
                headers: Iterable[str] = {k for item in value for k in item.keys()}
                ws.append(list(headers))
                for item in value:
                    ws.append([item.get(h, "") for h in headers])
            else:
                ws.append([key])
                for item in value:
                    ws.append([item])
        elif isinstance(value, dict):
            ws.append(["Field", "Value"])
            for k, v in value.items():
                ws.append([k, v])
        else:
            ws.append([key, value])

        # Simple column sizing
        for col in ws.columns:
            values = [str(c.value) if c.value is not None else "" for c in col]
            max_len = max(len(v) for v in values) if values else 0
            col_letter = get_column_letter(col[0].column)
            ws.column_dimensions[col_letter].width = max(10, min(50, max_len + 2))

    meta_ws = wb.create_sheet(title="Metadata")
    for k, v in metadata.items():
        meta_ws.append([k, v])
    meta_ws.append(["generated_on", metadata.get("date", date.today().isoformat())])

    wb.save(workbook_path)
    return str(workbook_path)


def _add_bullet_slide(prs: "Presentation", title: str, items: Dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    first = True
    for k, v in items.items():
        text = f"{k}: {v}"
        if first:
            tf.text = text
            first = False
        else:
            p = tf.add_paragraph()
            p.text = text


def generate_ppt_deck(result_data: JsonLike, metadata: Dict[str, Any], output_dir: Union[str, Path] = "Reports") -> str:
    """Generate a PowerPoint deck summarizing valuation results."""
    if Presentation is None:
        raise ImportError("python-pptx is required to generate PowerPoint decks")

    results = _load_json(result_data)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    case_name = metadata.get("case_name", "case")
    case_slug = _slugify(case_name) or "case"
    ppt_path = output_dir / f"deck_{case_slug}.pptx"

    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"Valuation Results - {case_name}"
    subtitle = slide.placeholders[1]
    subtitle.text = f"Date: {metadata.get('date', date.today().isoformat())}"

    # Assumptions slide
    assumptions = results.get("assumptions", {})
    if isinstance(assumptions, list) and assumptions and isinstance(assumptions[0], dict):
        assumptions = assumptions[0]
    if isinstance(assumptions, dict):
        _add_bullet_slide(prs, "Assumptions", assumptions)

    # Outputs slide
    outputs = results.get("outputs", {})
    if isinstance(outputs, list) and outputs and isinstance(outputs[0], dict):
        outputs = outputs[0]
    if isinstance(outputs, dict):
        _add_bullet_slide(prs, "Outputs", outputs)

    prs.save(ppt_path)
    return str(ppt_path)
